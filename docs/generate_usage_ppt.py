# -*- coding: utf-8 -*-
"""Generate the Quick_Sparam feature and usage PowerPoint.

The deck is built from live PyQt screenshots so the training material stays
aligned with the current UI.
"""

from __future__ import annotations

import math
import os
import sys
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = ROOT / "docs"
ASSET_DIR = DOCS_DIR / "ppt_assets"
OUTPUT_PPT = DOCS_DIR / "Quick_Sparam_功能介绍与使用指南.pptx"
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

FONT_FAMILY = "Microsoft YaHei UI"
TITLE = RGBColor(28, 39, 52)
BODY = RGBColor(45, 52, 61)
MUTED = RGBColor(98, 107, 118)
TEAL = RGBColor(23, 126, 137)
AMBER = RGBColor(229, 159, 46)
RED = RGBColor(218, 62, 67)
LIGHT_BG = RGBColor(248, 250, 252)
PANEL = RGBColor(238, 244, 246)


def ensure_paths() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def _qt_app():
    """Create a QApplication using the native Windows platform for CJK fonts."""
    from qtpy.QtWidgets import QApplication
    from qtpy.QtGui import QFont

    os.environ.pop("QT_QPA_PLATFORM", None)
    app = QApplication.instance() or QApplication(sys.argv)
    app.setFont(QFont(FONT_FAMILY, 9))
    return app


def _process_events(app, rounds: int = 8) -> None:
    for _ in range(rounds):
        app.processEvents()


def grab_widget(widget, path: Path, app, size: tuple[int, int] | None = None) -> None:
    if size:
        widget.resize(*size)
    widget.show()
    _process_events(app)
    pixmap = widget.grab()
    if not pixmap.save(str(path)):
        raise RuntimeError(f"截图保存失败: {path}")
    widget.close()
    _process_events(app, rounds=3)


def capture_ui_screenshots() -> dict[str, Path]:
    from qtpy.QtWidgets import QListWidgetItem

    from main_window import SParameterViewer_MainWin
    from QS_dialogs.cascade import CascadeDialog
    from QS_dialogs.freq_analysis import frequencyAnalysisDialog
    from QS_dialogs.port_management import PortManagementDialog
    from QS_dialogs.ripple import RippleFitDialog
    from QS_dialogs.se2diff import DiffConversionDialog
    from QS_dialogs.time_domain import TimeDomainDialog

    app = _qt_app()
    original_stdout = sys.stdout

    paths = {
        "main": ASSET_DIR / "main_window.png",
        "port_management": ASSET_DIR / "port_management.png",
        "diff": ASSET_DIR / "diff_conversion.png",
        "frequency": ASSET_DIR / "frequency_analysis.png",
        "cascade": ASSET_DIR / "cascade.png",
        "ripple": ASSET_DIR / "ripple.png",
        "time": ASSET_DIR / "time_domain.png",
    }

    main = SParameterViewer_MainWin()
    sys.stdout = original_stdout

    sample_files = [
        ROOT / "samples" / "StackupDemo1_test.s4p",
        ROOT / "samples" / "Twinax line-Spara1G.s4p",
        ROOT / "samples" / "connector_model.s4p",
    ]
    for sample in sample_files:
        main.file_list.addItem(str(sample))
    for i in range(main.file_list.count()):
        item = main.file_list.item(i)
        item.setSelected(i < 2)

    main.port1_input.setText("1 3")
    main.port2_input.setText("2 4")
    main.param_type_combo.setCurrentText("S参数")
    main.facet_combo.setCurrentText("幅度(dB)")
    main.mapping_combo.setCurrentText("一 一对应")
    main.freG_input.setText("10")
    main.same_plot_checkbox.setChecked(True)
    main.output_console.setPlainText(
        "演示流程:\n"
        "1. 打开一个或多个 .snp 文件\n"
        "2. 输入端口对并选择数据切面\n"
        "3. 点击绘图或进入分析/转换模块\n"
        "4. 在信息输出区查看处理日志"
    )
    grab_widget(main, paths["main"], app, (1500, 700))
    sys.stdout = original_stdout

    port_dialog = PortManagementDialog(None)
    grab_widget(port_dialog, paths["port_management"], app, (420, 300))

    diff_dialog = DiffConversionDialog(4, None)
    diff_dialog.partial_diff_radio.setChecked(True)
    diff_dialog.diff_lines_edit.setText("1 2")
    grab_widget(diff_dialog, paths["diff"], app)

    freq_dialog = frequencyAnalysisDialog({}, None)
    freq_dialog.analysis_checks["insertion_loss"].setChecked(True)
    freq_dialog.analysis_checks["return_loss"].setChecked(True)
    freq_dialog.analysis_checks["group_delay"].setChecked(True)
    freq_dialog.line_input.setText("1:4")
    freq_dialog.freG_input.setText("1 5 10 20")
    grab_widget(freq_dialog, paths["frequency"], app, (950, 520))

    cascade_dialog = CascadeDialog(None, ["connector_model.s4p", "board_model.s4p"])
    if cascade_dialog.table.rowCount() >= 2:
        cascade_dialog.table.item(0, 1).setText("1 2")
        cascade_dialog.table.item(0, 2).setText("3 4")
        cascade_dialog.table.item(1, 1).setText("3 4")
        cascade_dialog.table.item(1, 2).setText("1 2")
        cascade_dialog.set_pair_colors()
    grab_widget(cascade_dialog, paths["cascade"], app, (1200, 500))

    ripple_dialog = RippleFitDialog(None, [str(sample_files[0])])
    ripple_dialog.port1_input.setText("1")
    ripple_dialog.port2_input.setText("2")
    ripple_dialog.stop_freq_input.setText("20")
    ripple_dialog.fit_method.setCurrentText("IEEE_std_802.3-2022")
    grab_widget(ripple_dialog, paths["ripple"], app, (620, 330))

    td_dialog = TimeDomainDialog({}, None)
    td_dialog._port1_edit.setText("1")
    td_dialog._port2_edit.setText("1")
    td_dialog._port_list.addItem(QListWidgetItem("StackupDemo1_test.s4p  S1,1  [TDR]"))
    grab_widget(td_dialog, paths["time"], app, (900, 680))

    sys.stdout = original_stdout
    return paths


def font_path() -> str | None:
    candidates = [
        Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts" / "msyh.ttc",
        Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts" / "simhei.ttf",
        Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts" / "simsun.ttc",
    ]
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    return None


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    path = font_path()
    if path:
        try:
            return ImageFont.truetype(path, size=size, index=0)
        except Exception:
            pass
    return ImageFont.load_default()


def norm_rect(rect: tuple[float, float, float, float], width: int, height: int) -> tuple[int, int, int, int]:
    x1, y1, x2, y2 = rect
    return (int(x1 * width), int(y1 * height), int(x2 * width), int(y2 * height))


def norm_point(point: tuple[float, float], width: int, height: int) -> tuple[int, int]:
    return (int(point[0] * width), int(point[1] * height))


def draw_arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color, width: int = 5) -> None:
    draw.line([start, end], fill=color, width=width)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    length = max(16, width * 4)
    wing = math.pi / 7
    p1 = (
        end[0] - length * math.cos(angle - wing),
        end[1] - length * math.sin(angle - wing),
    )
    p2 = (
        end[0] - length * math.cos(angle + wing),
        end[1] - length * math.sin(angle + wing),
    )
    draw.polygon([end, p1, p2], fill=color)


def draw_label(
    draw: ImageDraw.ImageDraw,
    xy: tuple[int, int],
    number: int,
    text: str,
    font: ImageFont.ImageFont,
    small_font: ImageFont.ImageFont,
) -> tuple[int, int]:
    x, y = xy
    circle_r = 22
    line_gap = 8
    color = (218, 62, 67)
    shadow = (255, 255, 255, 235)
    text_bbox = draw.textbbox((0, 0), text, font=font)
    box_w = text_bbox[2] - text_bbox[0] + 28
    box_h = max(circle_r * 2, text_bbox[3] - text_bbox[1] + 24)
    box = (x, y, x + circle_r * 2 + line_gap + box_w, y + box_h)
    draw.rounded_rectangle(box, radius=14, fill=shadow, outline=color, width=3)
    draw.ellipse((x, y, x + circle_r * 2, y + circle_r * 2), fill=color)
    num_text = str(number)
    num_bbox = draw.textbbox((0, 0), num_text, font=small_font)
    draw.text(
        (x + circle_r - (num_bbox[2] - num_bbox[0]) / 2, y + circle_r - (num_bbox[3] - num_bbox[1]) / 2 - 2),
        num_text,
        fill=(255, 255, 255),
        font=small_font,
    )
    draw.text((x + circle_r * 2 + line_gap + 14, y + 10), text, fill=(28, 39, 52), font=font)
    return (x + circle_r, y + circle_r)


def annotate_image(
    src: Path,
    dst: Path,
    annotations: Iterable[dict],
) -> Path:
    image = Image.open(src).convert("RGBA")
    overlay = Image.new("RGBA", image.size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(overlay)
    w, h = image.size
    label_font = load_font(max(24, int(w * 0.016)))
    number_font = load_font(max(24, int(w * 0.018)))

    prepared = []
    for ann in annotations:
        rect = norm_rect(ann["rect"], w, h)
        x1, y1, x2, y2 = rect
        draw.rounded_rectangle(rect, radius=10, outline=(218, 62, 67, 255), width=max(4, w // 320))
        draw.rectangle((x1, y1, x2, y2), fill=(218, 62, 67, 24))
        prepared.append((ann, rect))

    for ann, rect in prepared:
        x1, y1, x2, y2 = rect
        label_xy = norm_point(ann["label_xy"], w, h)
        circle_center = draw_label(
            draw,
            label_xy,
            ann["num"],
            ann["text"],
            label_font,
            number_font,
        )
        target = ann.get("target")
        if target:
            target_xy = norm_point(target, w, h)
        else:
            target_xy = ((x1 + x2) // 2, (y1 + y2) // 2)
        draw_arrow(draw, circle_center, target_xy, (218, 62, 67, 255), width=max(4, w // 420))

    result = Image.alpha_composite(image, overlay).convert("RGB")
    result.save(dst, quality=95)
    return dst


def create_annotated_assets(paths: dict[str, Path]) -> dict[str, Path]:
    annotated = {}
    specs = {
        "main": [
            {"num": 1, "text": "文件操作", "rect": (0.011, 0.034, 0.135, 0.315), "label_xy": (0.145, 0.045)},
            {"num": 2, "text": "S参数操作", "rect": (0.011, 0.342, 0.135, 0.682), "label_xy": (0.145, 0.34)},
            {"num": 3, "text": "文件列表", "rect": (0.159, 0.034, 0.704, 0.642), "label_xy": (0.43, 0.16)},
            {"num": 4, "text": "绘图控制", "rect": (0.159, 0.679, 0.704, 0.977), "label_xy": (0.43, 0.74)},
            {"num": 5, "text": "信息输出", "rect": (0.718, 0.034, 0.994, 0.982), "label_xy": (0.70, 0.72)},
        ],
        "quick_plot": [
            {"num": 1, "text": "打开 .snp 文件", "rect": (0.017, 0.068, 0.129, 0.113), "label_xy": (0.145, 0.08)},
            {"num": 2, "text": "选中文件", "rect": (0.166, 0.060, 0.697, 0.626), "label_xy": (0.39, 0.10)},
            {"num": 3, "text": "输入端口对", "rect": (0.257, 0.704, 0.428, 0.767), "label_xy": (0.07, 0.77)},
            {"num": 4, "text": "选择参数与切面", "rect": (0.492, 0.745, 0.695, 0.872), "label_xy": (0.69, 0.60)},
            {"num": 5, "text": "点击绘图", "rect": (0.165, 0.829, 0.432, 0.904), "label_xy": (0.06, 0.88)},
            {"num": 6, "text": "查看日志", "rect": (0.724, 0.054, 0.986, 0.924), "label_xy": (0.69, 0.88)},
        ],
        "port_management": [
            {"num": 1, "text": "元数据", "rect": (0.05, 0.08, 0.95, 0.31), "label_xy": (0.08, 0.32)},
            {"num": 2, "text": "拓扑变换", "rect": (0.05, 0.34, 0.95, 0.58), "label_xy": (0.08, 0.59)},
            {"num": 3, "text": "阻抗变换", "rect": (0.05, 0.61, 0.95, 0.80), "label_xy": (0.08, 0.80)},
        ],
        "diff": [
            {"num": 1, "text": "选择线/端口逻辑", "rect": (0.03, 0.04, 0.97, 0.14), "label_xy": (0.05, 0.14)},
            {"num": 2, "text": "确认端口排布", "rect": (0.03, 0.16, 0.97, 0.49), "label_xy": (0.05, 0.48)},
            {"num": 3, "text": "设置差分阻抗", "rect": (0.03, 0.51, 0.97, 0.60), "label_xy": (0.05, 0.60)},
            {"num": 4, "text": "选择输出模式", "rect": (0.03, 0.76, 0.97, 0.86), "label_xy": (0.05, 0.86)},
        ],
        "frequency": [
            {"num": 1, "text": "勾选分析项目", "rect": (0.03, 0.05, 0.97, 0.23), "label_xy": (0.04, 0.24)},
            {"num": 2, "text": "指定端口排布", "rect": (0.03, 0.25, 0.97, 0.59), "label_xy": (0.04, 0.57)},
            {"num": 3, "text": "输入关注频点", "rect": (0.03, 0.61, 0.97, 0.78), "label_xy": (0.04, 0.76)},
            {"num": 4, "text": "绘图或导出 Excel", "rect": (0.03, 0.80, 0.97, 0.96), "label_xy": (0.46, 0.72)},
        ],
        "cascade": [
            {"num": 1, "text": "快速填端口", "rect": (0.02, 0.05, 0.98, 0.18), "label_xy": (0.04, 0.18)},
            {"num": 2, "text": "配置级联文件和端口", "rect": (0.02, 0.20, 0.98, 0.82), "label_xy": (0.04, 0.73)},
            {"num": 3, "text": "确认生成", "rect": (0.02, 0.84, 0.98, 0.96), "label_xy": (0.70, 0.70)},
        ],
        "ripple": [
            {"num": 1, "text": "端口对", "rect": (0.03, 0.06, 0.97, 0.35), "label_xy": (0.05, 0.36)},
            {"num": 2, "text": "频率范围", "rect": (0.03, 0.37, 0.97, 0.50), "label_xy": (0.05, 0.50)},
            {"num": 3, "text": "拟合方法", "rect": (0.03, 0.52, 0.97, 0.70), "label_xy": (0.05, 0.68)},
            {"num": 4, "text": "执行分析", "rect": (0.03, 0.80, 0.97, 0.95), "label_xy": (0.62, 0.68)},
        ],
        "time": [
            {"num": 1, "text": "波形类型", "rect": (0.02, 0.03, 0.98, 0.12), "label_xy": (0.04, 0.12)},
            {"num": 2, "text": "时域参数", "rect": (0.02, 0.14, 0.60, 0.36), "label_xy": (0.04, 0.36)},
            {"num": 3, "text": "兼容性检查", "rect": (0.62, 0.14, 0.98, 0.36), "label_xy": (0.58, 0.36)},
            {"num": 4, "text": "添加端口对", "rect": (0.02, 0.38, 0.98, 0.74), "label_xy": (0.04, 0.72)},
            {"num": 5, "text": "绘图控制", "rect": (0.02, 0.76, 0.98, 0.90), "label_xy": (0.63, 0.88)},
        ],
    }

    for key, source in paths.items():
        target = ASSET_DIR / f"{key}_annotated.png"
        spec_key = key if key != "main" else "main"
        annotate_image(source, target, specs[spec_key])
        annotated[key] = target

    quick_target = ASSET_DIR / "quick_plot_annotated.png"
    annotate_image(paths["main"], quick_target, specs["quick_plot"])
    annotated["quick_plot"] = quick_target
    return annotated


def set_fill(shape, color: RGBColor, transparency: int | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency


def add_text_box(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    color: RGBColor = BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullets(slide, items: Iterable[str], x: float, y: float, w: float, h: float, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT_FAMILY
        p.font.size = Pt(size)
        p.font.color.rgb = BODY
        p.space_after = Pt(5)
        p.level = 0


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    add_text_box(slide, title, 0.48, 0.22, 8.6, 0.38, size=24, color=TITLE, bold=True)
    if subtitle:
        add_text_box(slide, subtitle, 0.50, 0.63, 9.0, 0.28, size=12, color=MUTED)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
    set_fill(bar, TEAL)
    bar.line.fill.background()


def add_picture_fit(slide, image_path: Path, x: float, y: float, w: float, h: float, border: bool = True):
    with Image.open(image_path) as img:
        iw, ih = img.size
    max_w = Inches(w)
    max_h = Inches(h)
    ratio = min(max_w / iw, max_h / ih)
    pic_w = int(iw * ratio)
    pic_h = int(ih * ratio)
    left = Inches(x) + int((max_w - pic_w) / 2)
    top = Inches(y) + int((max_h - pic_h) / 2)
    pic = slide.shapes.add_picture(str(image_path), left, top, width=pic_w, height=pic_h)
    if border:
        line = pic.line
        line.color.rgb = RGBColor(205, 214, 221)
        line.width = Pt(0.75)
    return pic


def add_step(slide, num: int, title: str, text: str, x: float, y: float, w: float = 1.75) -> None:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.48), Inches(0.48))
    set_fill(circle, TEAL if num % 2 else AMBER)
    circle.line.fill.background()
    add_text_box(slide, str(num), x, y + 0.06, 0.48, 0.28, size=15, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, x + 0.58, y - 0.02, w, 0.28, size=15, color=TITLE, bold=True)
    add_text_box(slide, text, x + 0.58, y + 0.27, w, 0.55, size=11, color=MUTED)


def add_right_arrow(slide, x: float, y: float) -> None:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(0.55), Inches(0.26))
    set_fill(arrow, RGBColor(190, 205, 211))
    arrow.line.fill.background()


def deck_background(slide) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_BG


def create_deck(assets: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_text_box(slide, "Quick_Sparam", 0.75, 0.70, 5.6, 0.68, size=40, color=TITLE, bold=True)
    add_text_box(slide, "功能介绍与使用指南", 0.78, 1.45, 5.6, 0.45, size=23, color=TEAL, bold=True)
    add_text_box(
        slide,
        "面向 RF S 参数文件的查看、端口处理、级联、差分转换、频域/时域分析和结果导出",
        0.80,
        2.08,
        5.8,
        0.72,
        size=15,
        color=BODY,
    )
    add_bullets(
        slide,
        ["适用文件：Touchstone .snp", "核心对象：文件列表、端口对、数据切面、分析结果", "建议流程：先检查端口，再进行绘图或批量分析"],
        0.82,
        3.05,
        5.4,
        1.35,
        size=15,
    )
    add_picture_fit(slide, assets["main"], 6.15, 0.75, 6.6, 5.95)

    # 2. Workflow
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "典型操作流程", "从原始 Touchstone 文件到图形、转换文件和 Excel 报告。")
    steps = [
        ("导入文件", "打开一个或多个 .snp"),
        ("确认端口", "端口名、顺序、阻抗"),
        ("快速绘图", "输入端口对与数据切面"),
        ("分析转换", "频域、时域、级联、差分"),
        ("保存导出", "保存 S 参数或 Excel"),
    ]
    x = 0.55
    for idx, (title, text) in enumerate(steps, start=1):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.30), Inches(2.12), Inches(1.45))
        set_fill(card, RGBColor(255, 255, 255))
        card.line.color.rgb = RGBColor(202, 216, 222)
        add_step(slide, idx, title, text, x + 0.20, 2.58, w=1.18)
        if idx < len(steps):
            add_right_arrow(slide, x + 2.24, 2.90)
        x += 2.55
    add_bullets(
        slide,
        ["端口相关操作会生成新的网络对象，建议生成后立即查看文件信息或保存。", "批量分析前先选中文件，确保端口排布方式与模型一致。", "信息输出区会记录关键处理过程，便于复核和追溯。"],
        1.05,
        4.70,
        11.2,
        1.10,
        size=17,
    )

    # 3. Main overview
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "主界面模块总览", "红框和编号对应日常使用中最常触达的区域。")
    add_picture_fit(slide, assets["main"], 0.32, 1.0, 12.68, 6.1)

    # 4. Quick plot
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "快速绘制一条 S 参数曲线", "适合打开文件后快速检查插损、回损、相位、群延迟或 Z/Y 参数。")
    add_picture_fit(slide, assets["quick_plot"], 0.30, 1.0, 12.72, 6.1)

    # 5. Port management
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "端口处理", "先修正端口元数据，再做拓扑或阻抗变换，可以减少后续分析中的端口误配。")
    add_picture_fit(slide, assets["port_management"], 0.70, 1.15, 5.0, 4.95)
    add_bullets(
        slide,
        ["编辑端口名：为无端口名的模型补齐名称，便于选择器按名称填端口。", "修改参考阻抗：只更新 Z0 标注，不重新归一化矩阵。", "重排/合并/缩并：生成新的 S 参数对象，建议生成后查看文件信息。"],
        6.25,
        1.65,
        5.95,
        2.45,
        size=18,
    )
    add_text_box(slide, "推荐顺序：端口名 → 端口顺序 → 阻抗/缩并 → 保存结果", 6.30, 4.65, 5.6, 0.55, size=18, color=TEAL, bold=True)

    # 6. Differential conversion
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "差分转换", "支持按线逻辑或端口逻辑建立差分对，并选择只输出 SDD 或完整混合模式参数。")
    add_picture_fit(slide, assets["diff"], 0.45, 1.05, 7.05, 5.75)
    add_bullets(
        slide,
        ["按侧/按线排布决定差分线的默认配对方式。", "部分差分可只转换指定 line，适合混合单端与差分网络。", "完整混合模式会保留 SDD/SCD/SDC/SCC，便于后续串扰诊断。"],
        7.75,
        1.65,
        4.85,
        2.55,
        size=17,
    )

    # 7. Frequency analysis
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "频域批量分析", "面向插损、回损、串扰、PN skew、群延迟和 VTF 等指标的批量查看与导出。")
    add_picture_fit(slide, assets["frequency"], 0.35, 1.0, 8.0, 5.9)
    add_bullets(
        slide,
        ["先选中文件，再勾选需要的分析项目。", "端口排布与传输方向会影响 line/串扰计算结果。", "“指定 line/最差 line/横向比较”用于定位问题通道。", "“导出数据为 Excel”用于交付报告或二次处理。"],
        8.65,
        1.30,
        4.15,
        4.15,
        size=16,
    )

    # 8. Cascade
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "S 参数级联", "通过表格指定相邻网络的连接端口，相同颜色的单元格表示连接关系。")
    add_picture_fit(slide, assets["cascade"], 0.35, 1.05, 9.2, 5.55)
    add_bullets(
        slide,
        ["每一行选择一个待级联的 S 参数文件。", "右端口与下一行左端口数量必须一致。", "可用“所有端口/按边排布/按线排布”快速填表。"],
        9.85,
        1.65,
        3.0,
        2.6,
        size=16,
    )

    # 9. Ripple and time-domain
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "Ripple 与时域分析", "分别用于频域曲线纹波评估，以及 TDR/冲激/阶跃/脉冲响应查看。")
    add_picture_fit(slide, assets["ripple"], 0.35, 1.05, 6.0, 2.85)
    add_picture_fit(slide, assets["time"], 0.35, 4.12, 6.0, 2.85)
    add_bullets(
        slide,
        ["Ripple：输入端口对和频率范围，选择多项式、IEEE 802.3 或平滑函数拟合。", "时域：选择波形类型后设置上升沿、步长、点数和窗口，兼容性状态会给出风险提示。", "两类分析都建议先用主界面确认端口对，再进入对话框批量处理。"],
        6.75,
        1.40,
        5.85,
        3.85,
        size=17,
    )

    # 10. Delivery tips
    slide = prs.slides.add_slide(blank)
    deck_background(slide)
    add_header(slide, "交付与复核建议", "让结果可复现、可追溯、可复查。")
    add_bullets(
        slide,
        ["保存转换后的 S 参数文件时，保留能说明处理动作的后缀。", "Excel 结果用于报告正文，原始 .snp 与转换后 .snp 一起归档。", "关键信息可从“文件信息”“频点列表”和信息输出区复核。", "批量处理前确认端口排布图与实际封装/连接器拓扑一致。"],
        0.9,
        1.45,
        5.6,
        3.35,
        size=19,
    )
    add_picture_fit(slide, assets["main"], 6.45, 1.15, 6.35, 4.95)
    add_text_box(slide, f"文档生成脚本：{DOCS_DIR / 'generate_usage_ppt.py'}", 0.92, 5.80, 11.2, 0.32, size=12, color=MUTED)
    add_text_box(slide, f"输出文件：{OUTPUT_PPT.name}", 0.92, 6.18, 11.2, 0.32, size=12, color=MUTED)

    prs.save(OUTPUT_PPT)


def main() -> None:
    ensure_paths()
    raw_paths = capture_ui_screenshots()
    assets = create_annotated_assets(raw_paths)
    create_deck(assets)
    print(f"PPT 已生成: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
